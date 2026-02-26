#!/usr/bin/env python3
"""
Qwen-Image-Layered - Image Layer Decomposition
Gradio Interface for Dragonsuite

Decomposes images into multiple RGBA layers for advanced editing.
Uses Qwen/Qwen-Image-Layered (Apache 2.0 license).

Optimized for RTX 5070 Ti (16GB VRAM) with CPU offloading.
"""

import gc
import io
import os
import zipfile
from datetime import datetime
from pathlib import Path

import gradio as gr
import torch
from PIL import Image

# Configuration
OUTPUT_DIR = Path(os.path.expanduser("~/ai_generated/qwen-layered"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Global state
pipeline = None
device = "cuda" if torch.cuda.is_available() else "cpu"

# Dragon favicon as base64 SVG
DRAGON_FAVICON = """
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 64 64">
  <defs>
    <linearGradient id="dragonGrad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:#8B5CF6"/>
      <stop offset="100%" style="stop-color:#EC4899"/>
    </linearGradient>
  </defs>
  <circle cx="32" cy="32" r="30" fill="url(#dragonGrad)"/>
  <path d="M20 42 Q25 35 32 38 Q39 35 44 42 L42 38 Q38 32 32 35 Q26 32 22 38 Z" fill="#fff" opacity="0.9"/>
  <circle cx="24" cy="28" r="4" fill="#fff"/>
  <circle cx="40" cy="28" r="4" fill="#fff"/>
  <circle cx="25" cy="28" r="2" fill="#1a1a2e"/>
  <circle cx="41" cy="28" r="2" fill="#1a1a2e"/>
  <path d="M15 18 Q20 12 25 16" stroke="#fff" stroke-width="2" fill="none"/>
  <path d="M49 18 Q44 12 39 16" stroke="#fff" stroke-width="2" fill="none"/>
  <text x="32" y="52" text-anchor="middle" fill="#fff" font-size="8" font-weight="bold">LAYER</text>
</svg>
"""


def load_pipeline():
    """Load the Qwen-Image-Layered pipeline with memory optimization."""
    global pipeline

    if pipeline is not None:
        return pipeline

    print("Loading Qwen-Image-Layered pipeline...")
    print(f"Device: {device}")

    if torch.cuda.is_available():
        print(f"VRAM before load: {torch.cuda.memory_allocated() / 1024**3:.2f} GB")

    from diffusers import QwenImageLayeredPipeline

    pipeline = QwenImageLayeredPipeline.from_pretrained(
        "Qwen/Qwen-Image-Layered",
        torch_dtype=torch.bfloat16,
    )

    # Enable sequential CPU offloading for 16GB VRAM safety
    # Sequential is more aggressive - only one component on GPU at a time
    pipeline.enable_sequential_cpu_offload()
    pipeline.set_progress_bar_config(disable=None)

    if torch.cuda.is_available():
        print(f"VRAM after load: {torch.cuda.memory_allocated() / 1024**3:.2f} GB")

    print("Pipeline loaded successfully!")
    return pipeline


def decompose_image(
    image: Image.Image,
    layers: int,
    resolution: int,
    cfg_scale: float,
    steps: int,
    seed: int,
    use_random_seed: bool,
    progress=gr.Progress(track_tqdm=True),
) -> tuple:
    """Decompose an image into multiple layers."""
    if image is None:
        return None, None, None, "Please upload an image"

    try:
        # Load pipeline if needed
        pipe = load_pipeline()

        # Convert to RGBA if needed
        if image.mode != "RGBA":
            image = image.convert("RGBA")

        # Determine seed
        if use_random_seed:
            seed = torch.randint(0, 2**32 - 1, (1,)).item()

        generator = torch.Generator(device="cuda" if torch.cuda.is_available() else "cpu")
        generator.manual_seed(seed)

        print(f"Decomposing image into {layers} layers at {resolution}px (seed: {seed})")

        # Clear VRAM before inference
        if torch.cuda.is_available():
            torch.cuda.empty_cache()

        # Run pipeline
        with torch.inference_mode():
            output = pipe(
                image=image,
                layers=layers,
                resolution=resolution,
                true_cfg_scale=cfg_scale,
                num_inference_steps=steps,
                negative_prompt=" ",
                cfg_normalize=True,
                use_en_prompt=True,
                generator=generator,
                num_images_per_prompt=1,
            )

        layer_images = output.images[0]  # List of PIL RGBA images

        # Create output directory for this run
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        run_dir = OUTPUT_DIR / f"layers_{timestamp}"
        run_dir.mkdir(parents=True, exist_ok=True)

        # Save individual layers
        saved_paths = []
        for i, layer_img in enumerate(layer_images):
            layer_path = run_dir / f"layer_{i:02d}.png"
            layer_img.save(layer_path, "PNG")
            saved_paths.append(str(layer_path))
            print(f"Saved: {layer_path}")

        # Create ZIP file
        zip_path = run_dir / f"layers_{timestamp}.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, layer_img in enumerate(layer_images):
                # Save to bytes buffer
                buffer = io.BytesIO()
                layer_img.save(buffer, format="PNG")
                buffer.seek(0)
                zf.writestr(f"layer_{i:02d}.png", buffer.read())

        print(f"ZIP created: {zip_path}")

        # Create PPTX if python-pptx is available
        pptx_path = None
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Add title slide
            title_slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(title_slide_layout)

            # Add all layers to one slide (stacked)
            for i, layer_img in enumerate(layer_images):
                # Save layer temporarily
                temp_path = run_dir / f"_temp_layer_{i}.png"
                layer_img.save(temp_path, "PNG")

                # Add to slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                slide.shapes.add_picture(str(temp_path), left, top, width=width)

                # Clean up temp file
                temp_path.unlink()

            pptx_path = run_dir / f"layers_{timestamp}.pptx"
            prs.save(str(pptx_path))
            print(f"PPTX created: {pptx_path}")
        except ImportError:
            print("python-pptx not installed, skipping PPTX export")
        except Exception as e:
            print(f"PPTX export failed: {e}")

        status = f"Decomposed into {len(layer_images)} layers\n"
        status += f"Seed: {seed}\n"
        status += f"Saved to: {run_dir}\n"
        if pptx_path:
            status += f"PPTX: {pptx_path.name}"

        # Return gallery images, ZIP path, PPTX path, status
        return (
            layer_images,
            str(zip_path),
            str(pptx_path) if pptx_path else None,
            status,
        )

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        return None, None, None, f"Error: {str(e)}"


def clear_memory():
    """Clear GPU memory."""
    global pipeline
    if pipeline is not None:
        del pipeline
        pipeline = None
    gc.collect()
    if torch.cuda.is_available():
        torch.cuda.empty_cache()
        vram = torch.cuda.memory_allocated() / 1024**3
        return f"Memory cleared. VRAM: {vram:.2f} GB"
    return "Memory cleared"


def create_interface():
    """Create the Gradio interface."""

    with gr.Blocks(title="Qwen-Image-Layered") as demo:
        gr.HTML(f"""
        <div style="text-align: center; margin-bottom: 1rem;">
            <div style="display: inline-block; width: 48px; height: 48px; margin-bottom: 0.5rem;">
                {DRAGON_FAVICON}
            </div>
            <h1 style="margin: 0;">Qwen-Image-Layered</h1>
            <p style="color: #666; margin: 0.5rem 0;">Image Layer Decomposition for Advanced Editing</p>
        </div>
        """)

        with gr.Row():
            with gr.Column(scale=1):
                input_image = gr.Image(
                    label="Input Image",
                    type="pil",
                    sources=["upload", "clipboard"],
                )

                with gr.Row():
                    layers = gr.Slider(
                        minimum=2,
                        maximum=8,
                        value=4,
                        step=1,
                        label="Number of Layers",
                        info="How many layers to decompose into",
                    )

                    resolution = gr.Dropdown(
                        choices=[640, 1024],
                        value=640,
                        label="Resolution",
                        info="640 recommended, 1024 uses more VRAM",
                    )

                with gr.Row():
                    cfg_scale = gr.Slider(
                        minimum=1.0,
                        maximum=10.0,
                        value=4.0,
                        step=0.5,
                        label="CFG Scale",
                        info="Classifier-free guidance scale",
                    )

                    steps = gr.Slider(
                        minimum=20,
                        maximum=100,
                        value=50,
                        step=5,
                        label="Steps",
                        info="Inference steps (more = better quality)",
                    )

                with gr.Row():
                    seed = gr.Number(
                        value=777,
                        label="Seed",
                        precision=0,
                    )
                    use_random_seed = gr.Checkbox(
                        label="Random Seed",
                        value=True,
                    )

                with gr.Row():
                    decompose_btn = gr.Button("Decompose", variant="primary", size="lg")
                    clear_btn = gr.Button("Clear Memory", variant="secondary")

            with gr.Column(scale=2):
                gallery = gr.Gallery(
                    label="Layer Output",
                    columns=4,
                    height="auto",
                    object_fit="contain",
                )

                with gr.Row():
                    zip_file = gr.File(label="Download ZIP")
                    pptx_file = gr.File(label="Download PPTX")

                status_text = gr.Textbox(
                    label="Status",
                    interactive=False,
                    lines=4,
                )

        # Event handlers
        decompose_btn.click(
            fn=decompose_image,
            inputs=[
                input_image,
                layers,
                resolution,
                cfg_scale,
                steps,
                seed,
                use_random_seed,
            ],
            outputs=[gallery, zip_file, pptx_file, status_text],
        )

        clear_btn.click(
            fn=clear_memory,
            outputs=[status_text],
        )

        gr.HTML(f"""
        <div style="text-align: center; margin-top: 1rem; padding: 1rem; background: #faf5ff; border-radius: 8px;">
            <p style="margin: 0; color: #6b21a8;">
                <strong>Output directory:</strong> ~/ai_generated/qwen-layered/
            </p>
            <p style="margin: 0.5rem 0 0 0; color: #9333ea; font-size: 0.9em;">
                Uses ~14-16GB VRAM with CPU offloading. Close other GPU services before use.
            </p>
        </div>
        """)

        gr.Markdown("""
        ---
        **Model:** [Qwen/Qwen-Image-Layered](https://huggingface.co/Qwen/Qwen-Image-Layered) |
        **License:** Apache 2.0 |
        **Paper:** [arXiv:2512.15603](https://arxiv.org/abs/2512.15603)
        """)

    return demo


if __name__ == "__main__":
    print("Qwen-Image-Layered")
    print("==================")
    print(f"Device: {device}")
    print(f"Output directory: {OUTPUT_DIR}")

    demo = create_interface()
    demo.launch(
        server_name="0.0.0.0",
        server_port=8013,
        share=False,
        show_error=True,
        favicon_path="/srv/containers/edq/media/favicons/qwen-image.svg",
        allowed_paths=[str(OUTPUT_DIR)],
        theme=gr.themes.Soft(primary_hue="violet", secondary_hue="pink"),
    )
